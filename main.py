# -*- coding: utf-8 -*-
"""main.py

Updated version without unstructured[all-docs].
Uses PyPDFLoader + custom PPTXLoader (python-pptx).
"""

import os
import glob
from dotenv import load_dotenv

from langchain_community.document_loaders import DirectoryLoader, PyPDFLoader, TextLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.schema import Document
from pptx import Presentation

# class PPTXLoader:
#     def __init__(self, file_path: str):
#         self.file_path = file_path

#     def load(self):
#         prs = Presentation(self.file_path)
#         text = []
#         for slide in prs.slides:
#             for shape in slide.shapes:
#                 if hasattr(shape, "text"):
#                     text.append(shape.text)
#         return [Document(page_content="\n".join(text), metadata={"source": os.path.basename(self.file_path)})]


PATH = "/content/drive/MyDrive/waffles/"
DOCUMENTS_PATH = os.path.join(PATH, "data/ES")
DB_FAISS_PATH = os.path.join(PATH, 'db')

# ppt_files = glob.glob(os.path.join(DOCUMENTS_PATH, "**/*.pptx"), recursive=True)
# ppts = []
# for f in ppt_files:
#     loader = PPTXLoader(f)
#     ppts.extend(loader.load())

# pdf_loader = DirectoryLoader(
#     DOCUMENTS_PATH,
#     glob="**/*.pdf",
#     loader_cls=PyPDFLoader,
#     use_multithreading=True,
#     show_progress=True
# )
# pdfs = pdf_loader.load()

# documents = ppts + pdfs
# print(f"Loaded {len(documents)} documents.")

# text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
# texts = text_splitter.split_documents(documents)
# print(f"Split into {len(texts)} text chunks.")

# --- Embeddings + FAISS DB ---
embeddings = HuggingFaceEmbeddings(model_name='sentence-transformers/all-MiniLM-L6-v2', model_kwargs={'device': 'cpu'})
db = FAISS.from_documents(texts, embeddings)
db.save_local(DB_FAISS_PATH)

from langchain.prompts import PromptTemplate
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains import RetrievalQA
from langchain.chains import LLMChain

GOOGLE_API_KEY = os.getenv('GOOGLE_API_KEY')

prompt_template = """Use the following pieces of context to answer the question at the end. Elaborate the answer so that it is easy to understand.
Context: {context}
Question: {question}
Answer:
"""

prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])

llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash", google_api_key=GOOGLE_API_KEY)

non_rag_chain = LLMChain(
    llm=llm,
    prompt=PromptTemplate.from_template("{question}")
)

db = FAISS.load_local(DB_FAISS_PATH, embeddings, allow_dangerous_deserialization=True)

qa_chain = RetrievalQA.from_chain_type(
    llm=llm,
    chain_type='stuff',
    retriever=db.as_retriever(search_kwargs={'k': 2}),
    return_source_documents=True,
    chain_type_kwargs={'prompt': prompt}
)

print("Type 'exit' or 'quit' to stop.")

while True:
    query = input("> ")
    if query.lower() in ['exit', 'quit']:
        break

    try:
        res_rag = qa_chain.invoke({'query': query})
        answer_rag = res_rag["result"]
        source_docs = res_rag["source_documents"]

        print(answer_rag)
        if source_docs:
            print("\nSources:")
            for doc in source_docs:
                print(f"- {os.path.basename(doc.metadata.get('source', 'Unknown'))}")

    except Exception as e:
        print(f"\nError details: {e}")
